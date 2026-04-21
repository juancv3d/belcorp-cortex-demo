"""Generate Belcorp Cortex Code Masterclass PowerPoint - 8 slides."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Snowflake brand colors
SF_BLUE = RGBColor(0x29, 0xB5, 0xE8)
SF_MID_BLUE = RGBColor(0x11, 0x56, 0x7F)
SF_STAR_BLUE = RGBColor(0x71, 0xD3, 0xDC)
SF_ORANGE = RGBColor(0xFF, 0x9F, 0x36)
SF_GRAY = RGBColor(0x8A, 0x99, 0x9E)
SF_LIGHT = RGBColor(0xF0, 0xF8, 0xFF)
SF_BODY = RGBColor(0x5B, 0x5B, 0x5B)
SF_ICEBERG = RGBColor(0x00, 0x35, 0x45)
SF_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SF_DARK_TEAL = RGBColor(0x00, 0x4D, 0x61)
SF_DARKER = RGBColor(0x00, 0x2A, 0x36)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
LOGO_PATH = os.path.join(os.path.dirname(__file__), "snowflake_logo_white.png")


def add_bg(sl, color):
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = color


def rect(sl, l, t, w, h, color):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s


def rrect(sl, l, t, w, h, color):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s


def txt(sl, l, t, w, h, text, sz=18, color=SF_BODY, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    tb.text_frame.vertical_anchor = anchor
    p = tb.text_frame.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(sz)
    r.font.color.rgb = color; r.font.bold = bold; r.font.name = "Calibri"
    return tb


def bullets(sl, l, t, w, h, items, sz=14, color=SF_BODY, bc=SF_BLUE):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6); p.space_before = Pt(2)
        br = p.add_run(); br.text = "\u2022  "; br.font.size = Pt(sz)
        br.font.color.rgb = bc; br.font.bold = True; br.font.name = "Calibri"
        tr = p.add_run(); tr.text = item; tr.font.size = Pt(sz)
        tr.font.color.rgb = color; tr.font.name = "Calibri"
    return tb


def circle(sl, l, t, size, color, symbol, sz=22):
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    s.text_frame.word_wrap = False
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = s.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = symbol; r.font.size = Pt(sz)
    r.font.color.rgb = SF_WHITE; r.font.bold = True
    return s


def accent_line(sl):
    rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.06), SF_BLUE)


def logo(sl):
    if os.path.exists(LOGO_PATH):
        sl.shapes.add_picture(LOGO_PATH, Inches(10.5), Inches(0.4), Inches(2.2))


# ── Build ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITULO
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, SF_ICEBERG); accent_line(sl); logo(sl)

txt(sl, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5),
    "Cortex Code Masterclass", sz=44, color=SF_WHITE, bold=True)
txt(sl, Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.8),
    "Del User Story al Pull Request con AI", sz=24, color=SF_STAR_BLUE)
rect(sl, Inches(0.8), Inches(4.3), Inches(3.0), Inches(0.06), SF_BLUE)
txt(sl, Inches(0.8), Inches(5.2), Inches(5), Inches(0.5),
    "Belcorp  \u2022  Data Strategy Team", sz=16, color=SF_GRAY)
txt(sl, Inches(0.8), Inches(5.7), Inches(5), Inches(0.5),
    "Snowflake  \u2022  2025", sz=14, color=SF_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: QUE ES CORTEX CODE
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, SF_WHITE); accent_line(sl)

txt(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
    "Que es Cortex Code?", sz=32, color=SF_MID_BLUE, bold=True)

txt(sl, Inches(0.8), Inches(1.4), Inches(6.5), Inches(1.0),
    "AI coding agent nativo de Snowflake que vive en tu terminal.\n"
    "Conectado directamente a tus datos, sin extraer nada.",
    sz=18, color=SF_BODY)

bullets(sl, Inches(0.8), Inches(2.6), Inches(6.5), Inches(3.5), [
    "Entiende tu warehouse: tablas, columnas, tipos, relaciones",
    "Ejecuta SQL en tiempo real y valida resultados",
    "Lee y escribe codigo: SQL, Python, dbt, YAML",
    "Se conecta a GitHub, Jira, Confluence via MCP o CLI",
    "Planifica tareas, genera tests, crea Pull Requests",
], sz=16, color=SF_BODY)

rrect(sl, Inches(8.0), Inches(1.4), Inches(4.5), Inches(4.8), SF_ICEBERG)
txt(sl, Inches(8.4), Inches(1.8), Inches(3.7), Inches(0.5),
    "Diferenciador clave", sz=18, color=SF_BLUE, bold=True)
txt(sl, Inches(8.4), Inches(2.5), Inches(3.7), Inches(1.2),
    "No es un chatbot generico.\n\nEsta conectado a TUS datos\nen Snowflake y ejecuta\nacciones reales.",
    sz=15, color=SF_WHITE)
rect(sl, Inches(8.4), Inches(4.2), Inches(3.7), Inches(0.05), SF_BLUE)
txt(sl, Inches(8.4), Inches(4.5), Inches(3.7), Inches(1.2),
    "#TABLA \u2192 inyecta metadata\nSQL Execute \u2192 valida en vivo\ngh CLI \u2192 crea PRs directo",
    sz=14, color=SF_STAR_BLUE)

txt(sl, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
    "CLI  \u2022  Modelos Cortex  \u2022  Extensible con Skills, Agents y MCP",
    sz=13, color=SF_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: CAPACIDADES CLAVE (grid 2x3)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, SF_WHITE); accent_line(sl)

txt(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
    "Capacidades Clave", sz=32, color=SF_MID_BLUE, bold=True)

caps = [
    ("#", "#TABLE Syntax", "Inyecta metadata de tablas\nal prompt: columnas, tipos,\nPKs, sample rows"),
    ("S", "Custom Skills", "Archivos .md que estandarizan\nflujos y plantillas para\ntu equipo"),
    ("\u26A1", "SQL Execute", "Ejecuta y valida queries\nen tiempo real contra\nSnowflake"),
    ("\u21C4", "Multi-conexion", "DEV / QA / PRD desde\nun solo lugar via\nconnections.toml"),
    ("PR", "GitHub (gh CLI)", "Crea branches, commits\ny Pull Requests\nautomaticamente"),
    ("\u2713", "Task Tracking", "cortex ctx: planifica,\nrastrea progreso y\ngestiona pasos"),
]

for i, (icon, title, desc) in enumerate(caps):
    col, row = i % 3, i // 3
    x = Inches(0.6) + col * Inches(4.15)
    y = Inches(1.4) + row * Inches(2.8)
    rrect(sl, x, y, Inches(3.8), Inches(2.5), SF_LIGHT)
    rect(sl, x, y, Inches(3.8), Inches(0.06), SF_BLUE)
    circle(sl, x + Inches(0.25), y + Inches(0.3), Inches(0.6), SF_BLUE, icon, sz=18)
    txt(sl, x + Inches(1.05), y + Inches(0.3), Inches(2.5), Inches(0.5),
        title, sz=17, color=SF_MID_BLUE, bold=True)
    txt(sl, x + Inches(0.25), y + Inches(1.1), Inches(3.3), Inches(1.2),
        desc, sz=14, color=SF_BODY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: AGENDA
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, SF_WHITE); accent_line(sl)

txt(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
    "Agenda", sz=32, color=SF_MID_BLUE, bold=True)

# Left: Caso 1
rrect(sl, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2), SF_LIGHT)
rect(sl, Inches(0.8), Inches(1.5), Inches(5.6), Inches(0.08), SF_BLUE)
txt(sl, Inches(1.2), Inches(1.8), Inches(4.8), Inches(0.6),
    "Caso 1: Business-to-Engineering Copilot", sz=18, color=SF_MID_BLUE, bold=True)
bullets(sl, Inches(1.2), Inches(2.5), Inches(4.8), Inches(3.0), [
    "User Story \u2192 Refinamiento con AI",
    "Historia de Usuario \u2192 Historia Tecnica + SQL",
    "Generacion automatica de Test Cases",
    "Pull Request en GitHub",
], sz=15, color=SF_BODY)
rrect(sl, Inches(1.6), Inches(5.5), Inches(3.6), Inches(0.6), SF_BLUE)
txt(sl, Inches(1.6), Inches(5.5), Inches(3.6), Inches(0.6),
    "\u26A1  DEMO EN VIVO", sz=16, color=SF_WHITE, bold=True,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Right: Caso 2
rrect(sl, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.2), SF_LIGHT)
rect(sl, Inches(6.9), Inches(1.5), Inches(5.6), Inches(0.08), SF_ORANGE)
txt(sl, Inches(7.3), Inches(1.8), Inches(4.8), Inches(0.6),
    "Caso 2: Automated Testing", sz=18, color=SF_MID_BLUE, bold=True)
bullets(sl, Inches(7.3), Inches(2.5), Inches(4.8), Inches(3.0), [
    "Auto-generacion de tests desde historias",
    "Validacion cruzada DEV vs PRD",
    "Comparacion de datasets (data_diff)",
    "Evidencia + trazabilidad en PRs",
], sz=15, color=SF_BODY)
rrect(sl, Inches(7.7), Inches(5.5), Inches(3.6), Inches(0.6), SF_GRAY)
txt(sl, Inches(7.7), Inches(5.5), Inches(3.6), Inches(0.6),
    "CONCEPTUAL", sz=16, color=SF_WHITE, bold=True,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: CASO 1 - FLUJO (before demo)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, SF_WHITE); accent_line(sl)

txt(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
    "Caso 1: Business-to-Engineering Copilot", sz=28, color=SF_MID_BLUE, bold=True)
txt(sl, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
    "De requerimiento de negocio a codigo en produccion, con trazabilidad completa",
    sz=15, color=SF_GRAY)

steps = [
    ("1", "User Story", "Texto libre del negocio\nentra como prompt"),
    ("2", "Refinamiento AI", "Cortex Code refina con\ncriterios de aceptacion"),
    ("3", "Historia Tecnica\n+ SQL", "Genera queries validados\ncontra Snowflake"),
    ("4", "PR en GitHub", "Crea branch, commit y\nPull Request automatico"),
]

cw, ch = Inches(2.6), Inches(3.3)
sx = Inches(0.6)
gap = Inches(0.35)
cy = Inches(2.0)

for i, (num, title, desc) in enumerate(steps):
    x = sx + i * (cw + gap)
    rrect(sl, x, cy, cw, ch, SF_LIGHT)
    rect(sl, x, cy, cw, Inches(0.06), SF_BLUE)
    circle(sl, x + Inches(0.9), cy + Inches(0.3), Inches(0.7), SF_BLUE, num, sz=22)
    txt(sl, x + Inches(0.15), cy + Inches(1.2), Inches(2.3), Inches(0.8),
        title, sz=16, color=SF_MID_BLUE, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.15), cy + Inches(2.0), Inches(2.3), Inches(1.0),
        desc, sz=13, color=SF_BODY, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(sl, x + cw + Inches(0.05), cy + Inches(1.3), Inches(0.25), Inches(0.5),
            "\u2192", sz=28, color=SF_BLUE, bold=True, align=PP_ALIGN.CENTER)

txt(sl, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
    "Capacidades:  #TABLE syntax  \u2022  Custom Skills  \u2022  SQL Execute  \u2022  gh CLI  \u2022  cortex ctx",
    sz=13, color=SF_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: DEMO EN VIVO
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, SF_ICEBERG); accent_line(sl)

txt(sl, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.2),
    "Demo en Vivo", sz=48, color=SF_WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(sl, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.8),
    "User Story  \u2192  Refinamiento  \u2192  SQL + Tests  \u2192  Pull Request",
    sz=22, color=SF_STAR_BLUE, align=PP_ALIGN.CENTER)
rect(sl, Inches(5.2), Inches(4.5), Inches(3.0), Inches(0.05), SF_BLUE)
txt(sl, Inches(2.0), Inches(4.9), Inches(9.3), Inches(0.5),
    "BELCORP_ANALYTICS.COMERCIAL  \u2022  Cortex Code CLI",
    sz=15, color=SF_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: CASO 2 - FLUJO AUTOMATED TESTING
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, SF_WHITE); accent_line(sl)

txt(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
    "Caso 2: Automated Testing", sz=28, color=SF_MID_BLUE, bold=True)
txt(sl, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
    "Validacion automatizada entre ambientes con trazabilidad completa",
    sz=15, color=SF_GRAY)

test_steps = [
    ("\u2699", "Auto-generacion\nde Tests",
     "Cortex Code analiza la historia\ntecnica y genera casos de prueba\nSQL automaticamente:\nintegridad, NULLs, rangos, logica"),
    ("\u21C4", "Validacion Cruzada\nDEV vs PRD",
     "Multi-conexion Snowflake\n(connections.toml) ejecuta los\nmismos tests en ambientes\ndiferentes y compara resultados"),
    ("\u2262", "Comparacion de\nDatasets",
     "data_diff compara tablas entre\nesquemas: filas agregadas,\neliminadas, valores modificados,\ndiferencias de schema"),
    ("\u2713", "Evidencia +\nTrazabilidad",
     "Resultados se documentan\nautomaticamente en el Pull\nRequest de GitHub via gh CLI.\nHistoria \u2192 Tests \u2192 Evidencia \u2192 PR"),
]

for i, (icon, title, desc) in enumerate(test_steps):
    x = sx + i * (cw + gap)
    rrect(sl, x, cy, cw, Inches(3.8), SF_LIGHT)
    rect(sl, x, cy, cw, Inches(0.06), SF_ORANGE)
    circle(sl, x + Inches(0.9), cy + Inches(0.25), Inches(0.7), SF_ORANGE, icon, sz=18)
    txt(sl, x + Inches(0.15), cy + Inches(1.1), Inches(2.3), Inches(0.8),
        title, sz=15, color=SF_MID_BLUE, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.15), cy + Inches(2.0), Inches(2.3), Inches(1.6),
        desc, sz=12, color=SF_BODY, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(sl, x + cw + Inches(0.05), cy + Inches(1.3), Inches(0.25), Inches(0.5),
            "\u2192", sz=28, color=SF_ORANGE, bold=True, align=PP_ALIGN.CENTER)

txt(sl, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
    "Capacidades:  Multi-conexion  \u2022  data_diff  \u2022  SQL Execute  \u2022  gh CLI  \u2022  Custom Skills",
    sz=13, color=SF_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: CIERRE (Extensibilidad + Preguntas)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, SF_ICEBERG); accent_line(sl); logo(sl)

txt(sl, Inches(0.8), Inches(0.5), Inches(9), Inches(0.6),
    "Extensibilidad", sz=28, color=SF_WHITE, bold=True)
txt(sl, Inches(0.8), Inches(1.2), Inches(10), Inches(0.5),
    "Cortex Code se conecta a todo el ecosistema de desarrollo",
    sz=16, color=SF_STAR_BLUE)

ext = [
    ("MCP Servers", [
        "Jira: importar historias de usuario",
        "Confluence: leer documentacion tecnica",
        "GitHub: gestion completa de repos",
        "PostgreSQL, Filesystem, APIs custom",
    ]),
    ("Custom Skills + Agents", [
        "Skills .md para estandarizar flujos",
        "Agentes custom para tareas especializadas",
        "Hooks para validaciones automaticas",
        "Integracion con CI/CD pipelines",
    ]),
    ("Snowflake Nativo", [
        "Cortex Analyst: NL \u2192 SQL via semantic models",
        "Semantic Views: modelos de datos curados",
        "Multi-conexion: DEV / QA / PRD",
        "data_diff: comparacion de datasets",
    ]),
]

for i, (title, items) in enumerate(ext):
    x = Inches(0.8) + i * Inches(4.1)
    rrect(sl, x, Inches(2.0), Inches(3.7), Inches(3.6), SF_DARK_TEAL)
    txt(sl, x + Inches(0.3), Inches(2.2), Inches(3.1), Inches(0.5),
        title, sz=18, color=SF_BLUE, bold=True)
    bullets(sl, x + Inches(0.3), Inches(2.8), Inches(3.1), Inches(2.6),
            items, sz=13, color=SF_WHITE, bc=SF_STAR_BLUE)

rect(sl, Inches(0), Inches(6.2), SLIDE_W, Inches(1.3), SF_DARKER)
txt(sl, Inches(0.8), Inches(6.4), Inches(8), Inches(0.8),
    "\u00BFPreguntas?  \u2022  Siguiente paso: POC con sus datos reales",
    sz=20, color=SF_WHITE, anchor=MSO_ANCHOR.MIDDLE)
txt(sl, Inches(9.5), Inches(6.9), Inches(3.5), Inches(0.4),
    "\u00A9 2025 Snowflake Inc.", sz=11, color=SF_GRAY, align=PP_ALIGN.RIGHT)

# ── Save ─────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "Belcorp_Cortex_Code_Masterclass.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
